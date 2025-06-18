"""
Data export module for various output formats.
"""

import json
import logging
from pathlib import Path
from typing import Dict, Any, List
from datetime import datetime
import pandas as pd
import yaml
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.utils.dataframe import dataframe_to_rows
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


class DataExporter:
    """Export analyzed data to various formats."""
    
    def __init__(self):
        """Initialize data exporter."""
        self.logger = logging.getLogger(__name__)
        
    def export(self, data: Dict[str, Any], output_path: Path, formats: List[str], 
               metadata: Dict[str, Any] = None):
        """Export data to specified formats.
        
        Args:
            data: Data to export
            output_path: Base output path
            formats: List of formats ('json', 'excel', 'yaml', 'markdown', 'powerpoint')
            metadata: Additional metadata to include
        """
        base_name = output_path.stem
        output_dir = output_path.parent
        
        # Add metadata if provided
        export_data = {
            'metadata': metadata or {},
            'extraction_date': datetime.now().isoformat(),
            'data': data
        }
        
        for format_type in formats:
            try:
                if format_type == "json":
                    self._export_json(export_data, output_dir / f"{base_name}.json")
                elif format_type == "excel":
                    self._export_excel(export_data, output_dir / f"{base_name}.xlsx")
                elif format_type == "yaml":
                    self._export_yaml(export_data, output_dir / f"{base_name}.yaml")
                elif format_type == "markdown":
                    self._export_markdown(export_data, output_dir / f"{base_name}.md")
                elif format_type == "powerpoint":
                    self._export_powerpoint(export_data, output_dir / f"{base_name}.pptx")
                else:
                    self.logger.warning(f"Unknown format: {format_type}")
            except Exception as e:
                self.logger.error(f"Error exporting to {format_type}: {e}")
                raise
    
    def _export_json(self, data: Dict[str, Any], output_path: Path):
        """Export data as JSON."""
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        self.logger.info(f"Exported JSON: {output_path}")
    
    def _export_yaml(self, data: Dict[str, Any], output_path: Path):
        """Export data as YAML."""
        with open(output_path, 'w', encoding='utf-8') as f:
            yaml.dump(data, f, allow_unicode=True, default_flow_style=False)
        self.logger.info(f"Exported YAML: {output_path}")
    
    def _export_excel(self, data: Dict[str, Any], output_path: Path):
        """Export data as Excel."""
        wb = Workbook()
        ws = wb.active
        ws.title = "Knowledge Extraction"
        
        # Style configuration
        header_font = Font(bold=True, size=12)
        header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
        header_alignment = Alignment(horizontal="center", vertical="center")
        
        # Headers
        ws['A1'] = "Category"
        ws['B1'] = "Item"
        ws['A1'].font = header_font
        ws['A1'].fill = header_fill
        ws['A1'].alignment = header_alignment
        ws['B1'].font = header_font
        ws['B1'].fill = header_fill
        ws['B1'].alignment = header_alignment
        
        # Data
        row = 2
        for category, items in data['data'].items():
            for item in items:
                ws[f'A{row}'] = category
                ws[f'B{row}'] = item
                row += 1
        
        # Adjust column widths
        ws.column_dimensions['A'].width = 25
        ws.column_dimensions['B'].width = 80
        
        # Add metadata sheet
        ws2 = wb.create_sheet(title="Metadata")
        ws2['A1'] = "Property"
        ws2['B1'] = "Value"
        ws2['A1'].font = header_font
        ws2['B1'].font = header_font
        
        row = 2
        for key, value in data['metadata'].items():
            ws2[f'A{row}'] = key
            ws2[f'B{row}'] = str(value)
            row += 1
        
        ws2[f'A{row}'] = "Extraction Date"
        ws2[f'B{row}'] = data['extraction_date']
        
        wb.save(output_path)
        self.logger.info(f"Exported Excel: {output_path}")
    
    def _export_markdown(self, data: Dict[str, Any], output_path: Path):
        """Export data as Markdown."""
        lines = []
        
        # Header
        lines.append("# Knowledge Extraction Report")
        lines.append("")
        
        # Metadata
        if data['metadata']:
            lines.append("## Document Information")
            for key, value in data['metadata'].items():
                lines.append(f"- **{key}**: {value}")
            lines.append("")
        
        lines.append(f"**Extraction Date**: {data['extraction_date']}")
        lines.append("")
        
        # Content
        lines.append("## Extracted Knowledge")
        lines.append("")
        
        for category, items in data['data'].items():
            lines.append(f"### {category}")
            lines.append("")
            for item in items:
                lines.append(f"- {item}")
            lines.append("")
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(lines))
        
        self.logger.info(f"Exported Markdown: {output_path}")
    
    def _export_powerpoint(self, data: Dict[str, Any], output_path: Path):
        """Export data as PowerPoint."""
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Knowledge Extraction Report"
        subtitle.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        
        # Content slides
        for category, items in data['data'].items():
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = category
            
            tf = body_shape.text_frame
            tf.clear()
            
            for item in items[:5]:  # Limit to 5 items per slide
                p = tf.add_paragraph()
                p.text = item
                p.level = 0
                p.font.size = Pt(18)
        
        prs.save(output_path)
        self.logger.info(f"Exported PowerPoint: {output_path}")