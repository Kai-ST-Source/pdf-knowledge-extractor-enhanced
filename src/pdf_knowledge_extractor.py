#!/usr/bin/env python3
"""
PDF Knowledge Extractor - Enhanced Version
Extracts knowledge from PDF files with multiple extraction modes
"""

import os
import sys
import json
import argparse
import tempfile
import shutil
import __main__
import threading
import tkinter as tk
from tkinter import messagebox, ttk, filedialog
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Any, Optional, Tuple
import logging
from concurrent.futures import ThreadPoolExecutor, as_completed

import pandas as pd
from pdf2image import convert_from_path
from PIL import Image
import fitz  # PyMuPDF
from pptx import Presentation
import google.generativeai as genai
from tqdm import tqdm
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.utils.dataframe import dataframe_to_rows
import yaml
from pathvalidate import sanitize_filename

# Import our custom modules
from core.extractor import PDFExtractor
from core.analyzer import AIAnalyzer

# macOS specific imports
if sys.platform == "darwin":
    try:
        from Foundation import NSUserNotification, NSUserNotificationCenter
    except ImportError:
        NSUserNotification = None
        NSUserNotificationCenter = None

# Set environment variables for macOS Tkinter compatibility
os.environ['TK_SILENCE_DEPRECATION'] = '1'
if sys.platform == 'darwin':
    os.environ['PYTHON_CONFIGURE_OPTS'] = '--enable-framework'
    # Try to use system Tkinter
    try:
        tk._test()
    except Exception as e:
        print(f"Tkinter test failed: {e}")
        # Fallback to alternative display
        os.environ['DISPLAY'] = ':0'

class PDFKnowledgeExtractor:
    """Main class for extracting knowledge from PDF documents."""
    
    def __init__(self, config_path: str = "config.json"):
        """Initialize the extractor with configuration."""
        self.config = self._load_config(config_path)
        self.setup_logging()
        
        # Initialize Gemini client only if API key is provided and not in raw text mode
        self.gemini_client = None
        if self.config.get("gemini_api_key") and self.config.get("extraction_settings", {}).get("default_mode") != "raw_text_only":
            try:
                genai.configure(api_key=self.config["gemini_api_key"])
                self.gemini_client = genai.GenerativeModel(self.config["model_name"])
                logging.info("Gemini client initialized successfully")
            except Exception as e:
                logging.warning(f"Failed to initialize Gemini client: {e}")
                logging.info("Continuing without AI analysis - raw text extraction only")
        else:
            logging.info("No API key provided or raw text mode selected - AI analysis disabled")
        
        self.temp_dir = Path(tempfile.mkdtemp(prefix="pdf_extractor_"))
        self.results = []
        
        # Initialize core components
        self.extractor = PDFExtractor(self.temp_dir)
        
        # Initialize analyzer only if Gemini client is available
        self.analyzer = None
        if self.gemini_client:
            try:
                from core.analyzer import AIAnalyzer
                self.analyzer = AIAnalyzer(self.config)
                logging.info("AI Analyzer initialized successfully")
            except Exception as e:
                logging.warning(f"Failed to initialize AI Analyzer: {e}")
        else:
            logging.info("AI Analyzer disabled - no Gemini client available")
        
        self.output_dir = Path(self.config["output"]["output_directory"])
        self.output_dir.mkdir(exist_ok=True)
        
        logging.info("PDF Knowledge Extractor initialized successfully")
        
    def _load_config(self, config_path: str) -> Dict[str, Any]:
        """Load configuration from JSON file."""
        try:
            # Try multiple possible config file locations
            possible_paths = [
                config_path,
                Path(__file__).parent / config_path,
                Path(__file__).parent.parent / config_path,
                Path.home() / "Desktop" / "PDF knowledge extractor" / config_path,
                # For PyInstaller bundled app
                Path(sys._MEIPASS) / config_path if hasattr(sys, '_MEIPASS') else None
            ]
            
            # Filter out None values
            possible_paths = [p for p in possible_paths if p is not None]
            
            config_file = None
            for path in possible_paths:
                if Path(path).exists():
                    config_file = path
                    break
            
            if config_file is None:
                # Create default config if none exists
                default_config = {
                    "gemini_api_key": "",
                    "model_name": "gemini-1.5-flash",
                    "temperature": 0.3,
                    "max_tokens": 8192,
                    "log_level": "INFO",
                    "pdf_processing": {
                        "dpi": 200,
                        "max_images": 10
                    },
                    "output": {
                        "default_formats": ["json", "txt", "markdown"],
                        "output_directory": str(Path.home() / "Desktop" / "PDF knowledge extractor")
                    },
                    "extraction_settings": {
                        "default_mode": "raw_text_only",
                        "header_font_size_threshold": 14,
                        "footnote_position_threshold": 700,
                        "table_extraction_enabled": True,
                        "formatting_preservation": True,
                        "raw_extraction_formats": ["json", "txt", "markdown", "yaml"],
                        "include_font_info": True,
                        "include_position_info": True
                    }
                }
                
                # Save default config to user's desktop
                output_dir = Path.home() / "Desktop" / "PDF knowledge extractor"
                output_dir.mkdir(exist_ok=True)
                config_file = output_dir / "config.json"
                
                with open(config_file, 'w', encoding='utf-8') as f:
                    json.dump(default_config, f, ensure_ascii=False, indent=4)
                
                logging.info(f"Created default config file: {config_file}")
            
            with open(config_file, 'r', encoding='utf-8') as f:
                config = json.load(f)
                
            logging.info(f"Loaded config from: {config_file}")
            return config
            
        except Exception as e:
            logging.error(f"Error loading configuration: {e}")
            # Return minimal default config
            return {
                "gemini_api_key": "",
                "model_name": "gemini-1.5-flash",
                "temperature": 0.3,
                "max_tokens": 8192,
                "log_level": "INFO",
                "output": {
                    "default_formats": ["json", "txt", "markdown"],
                    "output_directory": str(Path.home() / "Desktop" / "PDF knowledge extractor")
                }
            }
    
    def setup_logging(self):
        """Setup logging configuration."""
        log_dir = Path.home() / "Desktop" / "PDF knowledge extractor"
        log_dir.mkdir(exist_ok=True)
        log_file = log_dir / "extraction.log"
        
        logging.basicConfig(
            level=getattr(logging, self.config.get("log_level", "INFO")),
            format='%(asctime)s - %(levelname)s - %(message)s',
            handlers=[
                logging.FileHandler(str(log_file)),
                logging.StreamHandler()
            ]
        )
        
    def process_file_detailed(self, file_path: Path, output_dir: Path, formats: List[str]) -> Dict[str, Any]:
        """Process a single file with detailed text extraction.
        
        Args:
            file_path: Path to PDF file
            output_dir: Output directory
            formats: List of output formats
            
        Returns:
            Dictionary of analysis results
        """
        logging.info(f"Processing with detailed extraction: {file_path}")
        
        # Send start notification
        self.send_notification("PDF Knowledge Extractor", f"詳細処理開始: {file_path.name}")
        
        try:
            # Extract detailed text information
            self.send_notification("PDF Knowledge Extractor", "PDFから詳細テキスト情報を抽出中...")
            detailed_text_info = self.extractor.extract_detailed_text(file_path)
            
            # Convert to images
            self.send_notification("PDF Knowledge Extractor", "PDFを画像に変換中...")
            images = self.extractor.extract_images(file_path)
            
            # Analyze with detailed information
            self.send_notification("PDF Knowledge Extractor", "AIで詳細知見を分析中...")
            # Analyze with AI if analyzer is available
            if self.analyzer:
                try:
                    results = self.analyzer.analyze_detailed(detailed_text_info, images)
                    logging.info("Detailed AI analysis completed successfully")
                except Exception as e:
                    logging.error(f"Error in detailed AI analysis: {e}")
                    results = {"error": f"AI analysis failed: {e}"}
            else:
                logging.info("Skipping AI analysis - analyzer not available")
                results = {"note": "AI analysis skipped - no API key or analyzer not available"}
            
            # Add metadata
            results['metadata'] = {
                'file_name': file_path.name,
                'file_path': str(file_path),
                'processed_at': datetime.now().isoformat(),
                'extraction_method': 'detailed',
                'pages_processed': len(detailed_text_info['pages']),
                'headers_found': len(detailed_text_info['headers']),
                'tables_found': len(detailed_text_info['tables']),
                'footnotes_found': len(detailed_text_info['footnotes'])
            }
            
            # Save results
            self.send_notification("PDF Knowledge Extractor", f"結果を保存中... ({', '.join(formats)})")
            output_path = output_dir / sanitize_filename(file_path.stem)
            self.save_results(results, output_path, formats)
            
            # Send completion notification
            self.send_notification("PDF Knowledge Extractor", f"詳細処理完了: {output_dir}")
            
            return results
            
        except Exception as e:
            logging.error(f"Error processing {file_path} with detailed extraction: {e}")
            self.send_notification("PDF Knowledge Extractor", f"エラーが発生しました: {str(e)}")
            raise
    
    def process_file(self, file_path: Path, output_dir: Path, formats: List[str]) -> Dict[str, Any]:
        """Process a single file with standard extraction.
        
        Args:
            file_path: Path to PDF file
            output_dir: Output directory
            formats: List of output formats
            
        Returns:
            Dictionary of analysis results
        """
        logging.info(f"Processing with standard extraction: {file_path}")
        
        # Send start notification
        self.send_notification("PDF Knowledge Extractor", f"処理開始: {file_path.name}")
        
        try:
            # Extract text and images
            self.send_notification("PDF Knowledge Extractor", "PDFからテキストを抽出中...")
            text, images = self.extractor.extract(file_path)
            
            # Analyze with AI
            self.send_notification("PDF Knowledge Extractor", "AIで知見を分析中...")
            # Analyze with AI if analyzer is available
            if self.analyzer:
                try:
                    results = self.analyzer.analyze(text, images)
                    logging.info("Standard AI analysis completed successfully")
                except Exception as e:
                    logging.error(f"Error in standard AI analysis: {e}")
                    results = {"error": f"AI analysis failed: {e}"}
            else:
                logging.info("Skipping AI analysis - analyzer not available")
                results = {"note": "AI analysis skipped - no API key or analyzer not available"}
            
            # Add metadata
            results['metadata'] = {
                'file_name': file_path.name,
                'file_path': str(file_path),
                'processed_at': datetime.now().isoformat(),
                'extraction_method': 'standard',
                'text_length': len(text)
            }
            
            # Save results
            self.send_notification("PDF Knowledge Extractor", f"結果を保存中... ({', '.join(formats)})")
            output_path = output_dir / sanitize_filename(file_path.stem)
            self.save_results(results, output_path, formats)
            
            # Send completion notification
            self.send_notification("PDF Knowledge Extractor", f"処理完了: {output_dir}")
            
            return results
            
        except Exception as e:
            logging.error(f"Error processing {file_path}: {e}")
            self.send_notification("PDF Knowledge Extractor", f"エラーが発生しました: {str(e)}")
            raise
    
    def process_file_raw_extraction(self, file_path: Path, output_dir: Path, formats: List[str]) -> Dict[str, Any]:
        """Process a single file with raw text extraction only (no AI analysis).
        
        Args:
            file_path: Path to PDF file
            output_dir: Output directory
            formats: List of output formats
            
        Returns:
            Dictionary of raw text data
        """
        logging.info(f"Processing with raw text extraction: {file_path}")
        
        # Send start notification
        self.send_notification("PDF Knowledge Extractor", f"生テキスト抽出開始: {file_path.name}")
        
        try:
            # Extract raw text information
            self.send_notification("PDF Knowledge Extractor", "PDFから生テキスト情報を抽出中...")
            raw_text_data = self.extractor.extract_raw_text_only(file_path)
            
            # Add metadata
            raw_text_data['metadata'] = {
                'file_name': file_path.name,
                'file_path': str(file_path),
                'processed_at': datetime.now().isoformat(),
                'extraction_method': 'raw_text_only',
                'total_pages': raw_text_data['total_pages'],
                'total_characters': len(raw_text_data['full_text']),
                'total_blocks': sum(len(page['blocks']) for page in raw_text_data['pages'])
            }
            
            # Save results
            self.send_notification("PDF Knowledge Extractor", f"結果を保存中... ({', '.join(formats)})")
            output_path = output_dir / sanitize_filename(file_path.stem)
            self.save_raw_results(raw_text_data, output_path, formats)
            
            # Send completion notification
            self.send_notification("PDF Knowledge Extractor", f"生テキスト抽出完了: {output_dir}")
            
            return raw_text_data
            
        except Exception as e:
            logging.error(f"Error processing {file_path} with raw extraction: {e}")
            self.send_notification("PDF Knowledge Extractor", f"エラーが発生しました: {str(e)}")
            raise
    
    def extract_text_from_pdf(self, pdf_path: Path) -> str:
        """Extract text from PDF using PyMuPDF."""
        try:
            doc = fitz.open(pdf_path)
            text = ""
            for page_num in range(len(doc)):
                page = doc[page_num]
                text += page.get_text()
            doc.close()
            return text
        except Exception as e:
            logging.error(f"Error extracting text from {pdf_path}: {e}")
            raise
    
    def convert_pdf_to_images(self, pdf_path: Path) -> List[Path]:
        """Convert PDF pages to images."""
        try:
            images = convert_from_path(pdf_path, dpi=200)
            image_paths = []
            
            for i, image in enumerate(images):
                image_path = self.temp_dir / f"page_{i+1}.png"
                image.save(image_path, 'PNG')
                image_paths.append(image_path)
                
            return image_paths
        except Exception as e:
            logging.error(f"Error converting PDF to images: {e}")
            raise
    
    def analyze_with_gemini(self, text: str, images: List[Path]) -> Dict[str, Any]:
        """Analyze document with Gemini AI."""
        try:
            prompt = self._build_analysis_prompt()
            
            # Prepare content for Gemini
            content = [prompt]
            
            # Add text if available
            if text.strip():
                content.append(f"Document Text:\n{text}")
            
            # Add images
            for image_path in images[:10]:  # Limit to first 10 images
                if image_path.exists():
                    content.append(Image.open(image_path))
            
            response = self.gemini_client.generate_content(content)
            return self._parse_gemini_response(response.text)
            
        except Exception as e:
            logging.error(f"Error analyzing with Gemini: {e}")
            raise
    
    def _build_analysis_prompt(self) -> str:
        """Build the analysis prompt for Gemini."""
        return """
あなたは専門的な文書分析エキスパートです。提供された文書から重要な知見を抽出し、以下の6つのカテゴリーに分類してください：

1. 概念・理論: 基本的な概念、理論、原理
2. 方法論・手順: 具体的な方法、手順、プロセス
3. 事例・ケーススタディ: 実例、事例研究、応用例
4. データ・数値: 統計データ、数値、測定結果
5. 注意点・リスク: 警告、注意事項、リスク要因
6. ベストプラクティス: 推奨事項、成功事例、最適解

各カテゴリーについて、以下の形式で回答してください：

1. 概念・理論
- 項目1: 説明
- 項目2: 説明

2. 方法論・手順
- 項目1: 説明
- 項目2: 説明

（以下同様）

重要なポイント：
- 各カテゴリーで最低1つ、最大5つの項目を抽出
- 簡潔で分かりやすい説明を心がける
- 文書の内容に忠実に
- 日本語で回答する
"""
    
    def _parse_gemini_response(self, response_text: str) -> Dict[str, Any]:
        """Parse Gemini AI response into structured data."""
        categories = {
            "概念・理論": [],
            "方法論・手順": [],
            "事例・ケーススタディ": [],
            "データ・数値": [],
            "注意点・リスク": [],
            "ベストプラクティス": []
        }
        
        current_category = None
        lines = response_text.split('\n')
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Check for category headers
            for i, category in enumerate(categories.keys(), 1):
                if f"{i}. {category}" in line or f"**{i}. {category}**" in line:
                    current_category = category
                    break
            
            # Extract items
            if current_category and line.startswith('- '):
                item = line[2:].strip()
                if item:
                    categories[current_category].append(item)
        
        return categories
    
    def save_results(self, results: Dict[str, Any], output_path: Path, formats: List[str]):
        """Save results in specified formats."""
        for format_type in formats:
            try:
                if format_type == "json":
                    self._save_json(results, output_path)
                elif format_type == "excel":
                    self._save_excel(results, output_path)
                elif format_type == "yaml":
                    self._save_yaml(results, output_path)
                elif format_type == "powerpoint":
                    self._save_powerpoint(results, output_path)
                elif format_type == "markdown":
                    self._save_markdown(results, output_path)
                else:
                    logging.warning(f"Unknown format: {format_type}")
            except Exception as e:
                logging.error(f"Error saving {format_type} format: {e}")
    
    def _save_json(self, results: Dict[str, Any], output_path: Path):
        """Save results as JSON."""
        json_path = output_path.with_suffix('.json')
        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump(results, f, ensure_ascii=False, indent=2)
        logging.info(f"Saved JSON: {json_path}")
    
    def _save_excel(self, results: Dict[str, Any], output_path: Path):
        """Save results as Excel file."""
        excel_path = output_path.with_suffix('.xlsx')
        
        wb = Workbook()
        
        # Remove default sheet
        wb.remove(wb.active)
        
        # Create sheets for each category
        for category, items in results.items():
            if category == 'metadata':
                continue
                
            ws = wb.create_sheet(title=category[:31])  # Excel sheet name limit
            
            # Add headers
            ws['A1'] = "項目"
            ws['B1'] = "説明"
            
            # Style headers
            for cell in ws[1]:
                cell.font = Font(bold=True)
                cell.fill = PatternFill(start_color="CCCCCC", end_color="CCCCCC", fill_type="solid")
            
            # Add data
            for i, item in enumerate(items, 2):
                ws[f'A{i}'] = f"項目{i-1}"
                ws[f'B{i}'] = item
            
            # Auto-adjust column widths
            for column in ws.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column_letter].width = adjusted_width
        
        wb.save(excel_path)
        logging.info(f"Saved Excel: {excel_path}")
    
    def _save_yaml(self, results: Dict[str, Any], output_path: Path):
        """Save results as YAML."""
        yaml_path = output_path.with_suffix('.yaml')
        with open(yaml_path, 'w', encoding='utf-8') as f:
            yaml.dump(results, f, default_flow_style=False, allow_unicode=True)
        logging.info(f"Saved YAML: {yaml_path}")
    
    def _save_powerpoint(self, results: Dict[str, Any], output_path: Path):
        """Save results as PowerPoint presentation."""
        pptx_path = output_path.with_suffix('.pptx')
        
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "PDF Knowledge Extraction Results"
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # Content slides
        for category, items in results.items():
            if category == 'metadata':
                continue
                
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            title.text = category
            
            content = slide.placeholders[1]
            content_text = ""
            for i, item in enumerate(items, 1):
                content_text += f"{i}. {item}\n"
            
            content.text = content_text
        
        prs.save(pptx_path)
        logging.info(f"Saved PowerPoint: {pptx_path}")
    
    def _save_markdown(self, results: Dict[str, Any], output_path: Path):
        """Save results as Markdown."""
        md_path = output_path.with_suffix('.md')
        
        with open(md_path, 'w', encoding='utf-8') as f:
            f.write("# PDF Knowledge Extraction Results\n\n")
            f.write(f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n")
            
            # Add metadata if available
            if 'metadata' in results:
                f.write("## Metadata\n\n")
                for key, value in results['metadata'].items():
                    f.write(f"- **{key}**: {value}\n")
                f.write("\n")
            
            # Add categories
            for category, items in results.items():
                if category == 'metadata':
                    continue
                    
                f.write(f"## {category}\n\n")
                for i, item in enumerate(items, 1):
                    f.write(f"{i}. {item}\n")
                f.write("\n")
        
        logging.info(f"Saved Markdown: {md_path}")
    
    def save_raw_results(self, raw_data: Dict[str, Any], output_path: Path, formats: List[str]):
        """Save raw text results in specified formats."""
        for format_type in formats:
            try:
                if format_type == "json":
                    self._save_raw_json(raw_data, output_path)
                elif format_type == "txt":
                    self._save_raw_txt(raw_data, output_path)
                elif format_type == "markdown":
                    self._save_raw_markdown(raw_data, output_path)
                elif format_type == "yaml":
                    self._save_raw_yaml(raw_data, output_path)
                else:
                    logging.warning(f"Unknown format for raw extraction: {format_type}")
            except Exception as e:
                logging.error(f"Error saving {format_type} format: {e}")
    
    def _save_raw_json(self, raw_data: Dict[str, Any], output_path: Path):
        """Save raw text data as JSON."""
        json_path = output_path.with_suffix('.raw.json')
        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump(raw_data, f, ensure_ascii=False, indent=2)
        logging.info(f"Saved raw JSON: {json_path}")
    
    def _save_raw_txt(self, raw_data: Dict[str, Any], output_path: Path):
        """Save raw text data as plain text."""
        txt_path = output_path.with_suffix('.raw.txt')
        with open(txt_path, 'w', encoding='utf-8') as f:
            f.write(f"=== PDF Raw Text Extraction ===\n")
            f.write(f"File: {raw_data['file_name']}\n")
            f.write(f"Pages: {raw_data['total_pages']}\n")
            f.write(f"Characters: {len(raw_data['full_text'])}\n")
            f.write(f"Extracted: {raw_data['extraction_timestamp']}\n")
            f.write("=" * 50 + "\n\n")
            
            for page in raw_data['pages']:
                f.write(f"--- PAGE {page['page_number']} ---\n")
                f.write(page['raw_text'])
                f.write("\n\n")
        
        logging.info(f"Saved raw TXT: {txt_path}")
    
    def _save_raw_markdown(self, raw_data: Dict[str, Any], output_path: Path):
        """Save raw text data as Markdown."""
        md_path = output_path.with_suffix('.raw.md')
        with open(md_path, 'w', encoding='utf-8') as f:
            f.write("# PDF Raw Text Extraction\n\n")
            f.write(f"**File:** {raw_data['file_name']}\n")
            f.write(f"**Pages:** {raw_data['total_pages']}\n")
            f.write(f"**Characters:** {len(raw_data['full_text'])}\n")
            f.write(f"**Extracted:** {raw_data['extraction_timestamp']}\n\n")
            
            for page in raw_data['pages']:
                f.write(f"## Page {page['page_number']}\n\n")
                f.write(f"**Text Length:** {page['text_length']} characters\n")
                f.write(f"**Blocks:** {len(page['blocks'])}\n\n")
                
                # Add raw text
                f.write("### Raw Text\n\n")
                f.write("```\n")
                f.write(page['raw_text'])
                f.write("\n```\n\n")
                
                # Add structured blocks
                if page['blocks']:
                    f.write("### Text Blocks\n\n")
                    for i, block in enumerate(page['blocks'], 1):
                        f.write(f"#### Block {i}\n\n")
                        f.write(f"**Text:** {block['text']}\n\n")
                        if block.get('font_info'):
                            font_info = block['font_info']
                            f.write(f"**Fonts:** {', '.join(font_info['fonts'])}\n")
                            f.write(f"**Sizes:** {', '.join(map(str, font_info['sizes']))}\n")
                            f.write(f"**Bold:** {font_info['is_bold']}\n")
                            f.write(f"**Italic:** {font_info['is_italic']}\n\n")
                
                f.write("---\n\n")
        
        logging.info(f"Saved raw Markdown: {md_path}")
    
    def _save_raw_yaml(self, raw_data: Dict[str, Any], output_path: Path):
        """Save raw text data as YAML."""
        yaml_path = output_path.with_suffix('.raw.yaml')
        with open(yaml_path, 'w', encoding='utf-8') as f:
            yaml.dump(raw_data, f, default_flow_style=False, allow_unicode=True)
        logging.info(f"Saved raw YAML: {yaml_path}")
    
    def send_notification(self, title: str, message: str):
        """Send system notification."""
        if sys.platform == "darwin" and NSUserNotification:
            notification = NSUserNotification.alloc().init()
            notification.setTitle_(title)
            notification.setInformativeText_(message)
            NSUserNotificationCenter.defaultUserNotificationCenter().deliverNotification_(notification)
        else:
            logging.info(f"Notification: {title} - {message}")
    
    def cleanup(self):
        """Clean up temporary files."""
        if self.temp_dir.exists():
            shutil.rmtree(self.temp_dir)
            logging.info(f"Cleaned up temporary directory: {self.temp_dir}")
        
        # Clean up extractor
        self.extractor.cleanup()

class PDFExtractorGUI:
    def __init__(self, extractor: PDFKnowledgeExtractor):
        """Initialize the GUI."""
        self.extractor = extractor
        self.output_dir = Path("/Users/hideki/Desktop/PDF knowledge extractor")
        self.formats = ["json", "excel", "markdown"]
        
        # Create main window with minimal setup
        self.root = tk.Tk()
        self.root.title("PDF Knowledge Extractor - Enhanced")
        self.root.geometry("600x500")
        
        # Simple window positioning
        self.root.update_idletasks()
        x = (self.root.winfo_screenwidth() // 2) - (600 // 2)
        y = (self.root.winfo_screenheight() // 2) - (500 // 2)
        self.root.geometry(f"600x500+{x}+{y}")
        
        # Setup UI
        self.setup_ui()
        
        # Protocol handler
        self.root.protocol("WM_DELETE_WINDOW", self.on_closing)
    
    def on_closing(self):
        """Handle window closing event."""
        self.root.quit()
        self.root.destroy()
        
    def setup_ui(self):
        # Title
        title_label = tk.Label(
            self.root, 
            text="PDF Knowledge Extractor",
            font=("Arial", 18, "bold"),
            bg="#f0f0f0",
            fg="#333333"
        )
        title_label.pack(pady=20)
        
        # Extraction mode selection frame
        mode_frame = tk.Frame(
            self.root,
            bg="#f0f0f0"
        )
        mode_frame.pack(pady=10)
        
        # Mode selection label
        mode_label = tk.Label(
            mode_frame,
            text="抽出モード:",
            font=("Arial", 12, "bold"),
            bg="#f0f0f0",
            fg="#333333"
        )
        mode_label.pack(side=tk.LEFT, padx=(0, 10))
        
        # Mode selection variable
        self.extraction_mode = tk.StringVar(value="raw_text_only")
        
        # Standard mode radio button
        standard_radio = tk.Radiobutton(
            mode_frame,
            text="標準抽出",
            variable=self.extraction_mode,
            value="standard",
            font=("Arial", 11),
            bg="#f0f0f0",
            fg="#333333"
        )
        standard_radio.pack(side=tk.LEFT, padx=5)
        
        # Detailed mode radio button
        detailed_radio = tk.Radiobutton(
            mode_frame,
            text="詳細抽出",
            variable=self.extraction_mode,
            value="detailed",
            font=("Arial", 11),
            bg="#f0f0f0",
            fg="#333333"
        )
        detailed_radio.pack(side=tk.LEFT, padx=5)
        
        # Raw text mode radio button
        raw_text_radio = tk.Radiobutton(
            mode_frame,
            text="生テキスト抽出",
            variable=self.extraction_mode,
            value="raw_text_only",
            font=("Arial", 11),
            bg="#f0f0f0",
            fg="#333333"
        )
        raw_text_radio.pack(side=tk.LEFT, padx=5)
        
        # Mode description
        mode_desc = tk.Label(
            self.root,
            text="生テキスト抽出: AI分析なしで大量のテキスト情報を抽出（推奨）\n詳細抽出: 見出し、表、脚注などの構造情報を活用\n標準抽出: 基本的なテキスト抽出とAI分析",
            font=("Arial", 9),
            bg="#f0f0f0",
            fg="#666666",
            wraplength=400,
            justify=tk.LEFT
        )
        mode_desc.pack(pady=5)
        
        # File selection frame
        self.select_frame = tk.Frame(
            self.root,
            bg="#ffffff",
            relief=tk.RAISED,
            borderwidth=2,
            width=400,
            height=150
        )
        self.select_frame.pack(pady=15, padx=50, fill=tk.BOTH, expand=True)
        self.select_frame.pack_propagate(False)
        
        # File selection label
        self.select_label = tk.Label(
            self.select_frame,
            text="📄\n\nクリックしてPDFファイルを選択\n\n複数ファイル選択可能",
            font=("Arial", 14),
            bg="#ffffff",
            fg="#666666",
            justify=tk.CENTER,
            cursor="hand2"
        )
        self.select_label.pack(expand=True)
        
        # Select button
        self.select_button = tk.Button(
            self.root,
            text="📁 PDFファイルを選択",
            font=("Arial", 12, "bold"),
            bg="#4CAF50",
            fg="white",
            relief=tk.FLAT,
            padx=20,
            pady=10,
            cursor="hand2",
            command=self.select_files
        )
        self.select_button.pack(pady=10)
        
        # Progress bar
        self.progress_var = tk.StringVar(value="待機中...")
        self.progress_label = tk.Label(
            self.root,
            textvariable=self.progress_var,
            font=("Arial", 10),
            bg="#f0f0f0",
            fg="#333333"
        )
        self.progress_label.pack(pady=5)
        
        self.progress_bar = ttk.Progressbar(
            self.root,
            mode='indeterminate',
            length=300
        )
        self.progress_bar.pack(pady=5)
        
        # Output location label
        output_label = tk.Label(
            self.root,
            text=f"出力先: {self.output_dir}",
            font=("Arial", 9),
            bg="#f0f0f0",
            fg="#666666"
        )
        output_label.pack(pady=5)
        
        # Bind click for file selection
        self.select_frame.bind("<Button-1>", self.select_files)
        self.select_label.bind("<Button-1>", self.select_files)
        
        # Change background on hover
        self.select_frame.bind("<Enter>", lambda e: self.select_frame.configure(bg="#f8f8f8"))
        self.select_frame.bind("<Leave>", lambda e: self.select_frame.configure(bg="#ffffff"))
        
    def select_files(self, event=None):
        files = filedialog.askopenfilenames(
            title="PDFファイルを選択",
            filetypes=[("PDF files", "*.pdf"), ("All files", "*.*")]
        )
        if files:
            self.process_files(list(files))
    
    def process_files(self, files):
        def worker():
            try:
                self.root.after(0, lambda: self.progress_var.set("処理中..."))
                self.root.after(0, lambda: self.progress_bar.start())
                
                for file_path in files:
                    file_path = Path(file_path)
                    self.root.after(0, lambda f=file_path.name: self.progress_var.set(f"処理中: {f}"))
                    
                    # Choose processing method based on selected mode
                    if self.extraction_mode.get() == "detailed":
                        results = self.extractor.process_file_detailed(file_path, self.output_dir, self.formats)
                    elif self.extraction_mode.get() == "raw_text_only":
                        results = self.extractor.process_file_raw_extraction(file_path, self.output_dir, self.formats)
                    else:
                        results = self.extractor.process_file(file_path, self.output_dir, self.formats)
                    
                self.root.after(0, lambda: self.progress_bar.stop())
                self.root.after(0, lambda: self.progress_var.set("処理完了！"))
                
                # Show completion message
                mode_text = {
                    "detailed": "詳細抽出",
                    "raw_text_only": "生テキスト抽出",
                    "standard": "標準抽出"
                }.get(self.extraction_mode.get(), "標準抽出")
                
                self.root.after(0, lambda: messagebox.showinfo(
                    "完了", 
                    f"処理が完了しました！\n\n抽出モード: {mode_text}\n出力先: {self.output_dir}\n\n生成ファイル: JSON, TXT, Markdown, YAML"
                ))
                
                self.root.after(2000, lambda: self.progress_var.set("待機中..."))
                
            except Exception as e:
                self.root.after(0, lambda: self.progress_bar.stop())
                self.root.after(0, lambda: self.progress_var.set("エラーが発生しました"))
                self.root.after(0, lambda: messagebox.showerror("エラー", f"処理中にエラーが発生しました:\n{str(e)}"))
                self.root.after(2000, lambda: self.progress_var.set("待機中..."))
        
        # Run processing in separate thread
        thread = threading.Thread(target=worker, daemon=True)
        thread.start()
    
    def run(self):
        self.root.mainloop()

def main():
    """Main entry point."""
    
    # Add debug logging at startup
    debug_log_path = "/Users/hideki/Desktop/PDF knowledge extractor/debug.log"
    Path(debug_log_path).parent.mkdir(exist_ok=True)
    
    with open(debug_log_path, "a", encoding="utf-8") as f:
        f.write(f"\n=== App Started: {datetime.now()} ===\n")
        f.write(f"Command line args: {sys.argv}\n")
        f.write(f"Working directory: {os.getcwd()}\n")
        f.write(f"Python version: {sys.version}\n")
        f.write(f"Platform: {sys.platform}\n")
    
    parser = argparse.ArgumentParser(description="Extract knowledge from PDF documents")
    parser.add_argument("input", nargs="?", help="Input PDF file or directory")
    parser.add_argument("-o", "--output", default="/Users/hideki/Desktop/PDF knowledge extractor", help="Output directory")
    parser.add_argument("-f", "--format", nargs="+", default=["json", "excel", "markdown"],
                       choices=["json", "excel", "yaml", "powerpoint", "markdown"],
                       help="Output formats")
    
    # Get the directory where the script is located
    # When frozen with PyInstaller, sys._MEIPASS contains the bundled directory
    if getattr(sys, 'frozen', False) and hasattr(sys, '_MEIPASS'):
        script_dir = Path(sys._MEIPASS)
    else:
        script_dir = Path(__file__).parent if hasattr(__main__, "__file__") else Path(sys.argv[0]).parent
    
    default_config = script_dir / "config.json"
    parser.add_argument("-c", "--config", default=str(default_config), help="Configuration file")
    
    args = parser.parse_args()
    
    # Log parsed arguments
    with open(debug_log_path, "a", encoding="utf-8") as f:
        f.write(f"Parsed input: {args.input}\n")
        f.write(f"Output dir: {args.output}\n")
        f.write(f"Formats: {args.format}\n")
        f.write(f"Config path: {args.config}\n")
        f.write(f"Config file exists: {Path(args.config).exists()}\n")
        f.write(f"Is frozen: {getattr(sys, 'frozen', False)}\n")
    
    # Initialize extractor
    try:
        with open(debug_log_path, "a", encoding="utf-8") as f:
            f.write("Initializing PDFKnowledgeExtractor...\n")
        extractor = PDFKnowledgeExtractor(args.config)
        with open(debug_log_path, "a", encoding="utf-8") as f:
            f.write("PDFKnowledgeExtractor initialized successfully\n")
    except Exception as e:
        with open(debug_log_path, "a", encoding="utf-8") as f:
            f.write(f"Error initializing PDFKnowledgeExtractor: {e}\n")
            import traceback
            f.write(f"Traceback: {traceback.format_exc()}\n")
        
        # Show error dialog
        try:
            import subprocess
            subprocess.run(['osascript', '-e', f'display dialog "初期化エラー: {str(e)}" buttons {{"OK"}} default button "OK" with icon stop'])
        except:
            pass
        return 1
    
    try:
        output_dir = Path(args.output)
        output_dir.mkdir(exist_ok=True)
        
        # Handle case where no input file is provided (e.g., from Finder right-click)
        if not args.input:
            with open(debug_log_path, "a", encoding="utf-8") as f:
                f.write("No input file provided - launching GUI mode\n")
            
            # For macOS PyInstaller apps, we need to handle GUI differently
            if sys.platform == "darwin" and getattr(sys, 'frozen', False):
                with open(debug_log_path, "a", encoding="utf-8") as f:
                    f.write("Running on frozen macOS app - using file dialog directly\n")
                
                # Use file dialog directly
                extractor.send_notification("PDF Knowledge Extractor", "PDFファイルを選択してください")
                
                import subprocess
                try:
                    # Use multiple file selection
                    result = subprocess.run([
                        'osascript', '-e',
                        'tell application "System Events" to return POSIX path of (choose file with prompt "PDFファイルを選択してください:" of type {"pdf"} with multiple selections allowed)'
                    ], capture_output=True, text=True, timeout=60)
                    
                    if result.returncode == 0 and result.stdout.strip():
                        # Parse selected files (AppleScript returns comma-separated list)
                        selected_files = result.stdout.strip().split(", ")
                        with open(debug_log_path, "a", encoding="utf-8") as f:
                            f.write(f"User selected files: {selected_files}\n")
                        
                        # Process each file
                        for file_path in selected_files:
                            file_path = Path(file_path.strip())
                            if file_path.exists() and file_path.suffix.lower() == '.pdf':
                                try:
                                    results = extractor.process_file(file_path, output_dir, args.format)
                                    extractor.send_notification("PDF Knowledge Extractor", f"処理完了: {file_path.name}")
                                except Exception as e:
                                    logging.error(f"Failed to process {file_path}: {e}")
                                    extractor.send_notification("PDF Knowledge Extractor", f"処理エラー: {file_path.name}")
                        
                        extractor.send_notification("PDF Knowledge Extractor", f"全ての処理が完了しました。\n出力先: {output_dir}")
                        return 0
                    else:
                        with open(debug_log_path, "a", encoding="utf-8") as f:
                            f.write("User cancelled file selection\n")
                        extractor.send_notification("PDF Knowledge Extractor", "ファイル選択がキャンセルされました")
                        return 0
                except subprocess.TimeoutExpired:
                    with open(debug_log_path, "a", encoding="utf-8") as f:
                        f.write("File dialog timed out\n")
                    extractor.send_notification("PDF Knowledge Extractor", "ファイル選択がタイムアウトしました")
                    return 0
                except Exception as e2:
                    with open(debug_log_path, "a", encoding="utf-8") as f:
                        f.write(f"File dialog error: {e2}\n")
                    extractor.send_notification("PDF Knowledge Extractor", f"ファイル選択エラー: {e2}")
                    return 1
            else:
                # Try to launch GUI mode for non-frozen apps
                try:
                    with open(debug_log_path, "a", encoding="utf-8") as f:
                        f.write("Creating GUI instance...\n")
                    gui = PDFExtractorGUI(extractor)
                    with open(debug_log_path, "a", encoding="utf-8") as f:
                        f.write("GUI instance created, starting mainloop...\n")
                    gui.run()
                    with open(debug_log_path, "a", encoding="utf-8") as f:
                        f.write("GUI mainloop ended\n")
                    return 0
                except Exception as e:
                    with open(debug_log_path, "a", encoding="utf-8") as f:
                        f.write(f"GUI launch error: {e}\n")
                        import traceback
                        f.write(f"Traceback: {traceback.format_exc()}\n")
                    
                    # Show error and exit
                    try:
                        import subprocess
                        subprocess.run(['osascript', '-e', f'display dialog "GUI起動エラー: {str(e)}" buttons {{"OK"}} default button "OK" with icon stop'])
                    except:
                        pass
                    return 1
        
        # Process provided input path
        input_path = Path(args.input)
        
        with open(debug_log_path, "a", encoding="utf-8") as f:
            f.write(f"Processing input path: {input_path}\n")
            f.write(f"Path exists: {input_path.exists()}\n")
            f.write(f"Is file: {input_path.is_file()}\n")
        
        if input_path.is_file():
            # Process single file
            results = extractor.process_file(input_path, output_dir, args.format)
            extractor.send_notification("PDF Knowledge Extractor", f"Successfully processed {input_path.name}")
            
        elif input_path.is_dir():
            # Process directory
            pdf_files = list(input_path.glob("*.pdf"))
            if not pdf_files:
                logging.warning(f"No PDF files found in {input_path}")
                return
            
            for pdf_file in tqdm(pdf_files, desc="Processing PDFs"):
                try:
                    extractor.process_file(pdf_file, output_dir, args.format)
                except Exception as e:
                    logging.error(f"Failed to process {pdf_file}: {e}")
                    continue
            
            extractor.send_notification("PDF Knowledge Extractor", f"Processed {len(pdf_files)} files")
        
        else:
            with open(debug_log_path, "a", encoding="utf-8") as f:
                f.write(f"Input path does not exist: {input_path}\n")
            logging.error(f"Input path does not exist: {input_path}")
            extractor.send_notification("PDF Knowledge Extractor", f"ファイルが見つかりません: {input_path}")
            return
            
    except Exception as e:
        logging.error(f"Application error: {e}")
        extractor.send_notification("PDF Knowledge Extractor", f"Error: {str(e)}")
        return 1
    
    finally:
        extractor.cleanup()
    
    return 0

if __name__ == "__main__":
    sys.exit(main())